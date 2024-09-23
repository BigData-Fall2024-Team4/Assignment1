import os
import pandas as pd
import zipfile
import json
import io
from azure.storage.blob import BlobServiceClient
from dotenv import load_dotenv
import PyPDF2
from docx import Document
from pptx import Presentation

# Load environment variables from .env file
load_dotenv()

# Set up Azure Blob Service
blob_service_client = BlobServiceClient.from_connection_string(os.getenv("AZURE_CONNECTION_STRING"))
container_name = os.getenv("AZURE_CONTAINER_NAME")

# Supported file formats
supported_formats = ['xlsx', 'txt', 'pdf', 'csv', 'docx', 'png', 'jpg', 'pptx', 'mp3', 'jsonID', 'pdb', 'py', 'zip']

def get_file_content(file_name):
    try:
        # Get the file extension
        file_extension = file_name.split('.')[-1].lower()

        # Check if the file format is supported
        if file_extension not in supported_formats:
            return "Error: File format not supported"

        # Get the blob client and download the file content
        blob_client = blob_service_client.get_blob_client(container=container_name, blob=file_name)
        file_content = blob_client.download_blob().readall()

        # Handle file content based on file extension
        if file_extension == 'txt':
            return file_content.decode('utf-8')
        elif file_extension == 'pdf':
            return convert_pdf_to_text(file_content)
        elif file_extension == 'xlsx':
            return convert_excel_to_text(file_content)
        elif file_extension == 'csv':
            return file_content.decode('utf-8')
        elif file_extension == 'docx':
            return convert_docx_to_text(file_content)
        elif file_extension in ['png', 'jpg']:
            return "Image file detected"
        elif file_extension == 'pptx':
            return convert_pptx_to_text(file_content)
        elif file_extension == 'mp3':
            return "Audio file detected"
        elif file_extension == 'jsonID':
            return json.loads(file_content.decode('utf-8'))
        elif file_extension == 'pdb':
            return "PDB file detected"
        elif file_extension == 'py':
            return file_content.decode('utf-8')
        elif file_extension == 'zip':
            return handle_zip_file(file_content)
        else:
            return "Error: Unsupported file format"

    except Exception as e:
        return f"Error: {str(e)}"

def convert_pdf_to_text(file_content):
    # Read PDF content
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text()
        return text if text else "PDF contains no readable text"
    except Exception as e:
        return f"Error processing PDF: {str(e)}"

def convert_excel_to_text(file_content):
    try:
        # Read Excel content using pandas
        excel_file = pd.read_excel(io.BytesIO(file_content))
        return excel_file.to_string()
    except Exception as e:
        return f"Error processing Excel file: {str(e)}"

def convert_docx_to_text(file_content):
    try:
        # Read DOCX content
        doc = Document(io.BytesIO(file_content))
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error processing DOCX file: {str(e)}"

def convert_pptx_to_text(file_content):
    try:
        # Read PPTX content
        prs = Presentation(io.BytesIO(file_content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return '\n'.join(text_runs) if text_runs else "PPTX contains no readable text"
    except Exception as e:
        return f"Error processing PPTX file: {str(e)}"

def handle_zip_file(file_content):
    try:
        # Handle ZIP file contents
        with zipfile.ZipFile(io.BytesIO(file_content)) as z:
            zip_content_list = z.namelist()
            return f"ZIP contains: {', '.join(zip_content_list)}"
    except Exception as e:
        return f"Error processing ZIP file: {str(e)}"
